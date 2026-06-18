import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

def set_run_lang(run, lang='zh-CN'):
    """Set language on a text run to prevent PowerPoint spell-check false positives"""
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('w:lang'), lang)

def set_textframe_lang(tf, lang='zh-CN'):
    """Set language on all runs in a text frame"""
    for para in tf.paragraphs:
        for run in para.runs:
            set_run_lang(run, lang)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

IMG_DIR = os.path.join(os.path.dirname(__file__), "img")
DARK = RGBColor(0x3A, 0x3F, 0x47)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_GRAY = RGBColor(0x55, 0x55, 0x55)
ACCENT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
TABLE_HEADER = RGBColor(0xF0, 0xF0, 0xF0)

# ===== 全局字体设置 =====
FONT_CN = '微软雅黑'
FONT_EN = 'Arial'

def set_font_on_run(run, cn=FONT_CN, en=FONT_EN):
    """Set both Chinese and English fonts on a text run"""
    run.font.name = cn
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    if latin is None:
        latin = etree.SubElement(rPr, qn('a:latin'))
    latin.set('typeface', en)
    ea_exist = rPr.find(qn('a:ea'))
    if ea_exist is None:
        ea_exist = etree.SubElement(rPr, qn('a:ea'))
    ea_exist.set('typeface', cn)
    sym = rPr.find(qn('a:sym'))
    if sym is None:
        sym = etree.SubElement(rPr, qn('a:sym'))
    sym.set('typeface', en)

def set_font_on_textframe(tf, cn=FONT_CN, en=FONT_EN):
    """Set font on all runs in a text frame"""
    for para in tf.paragraphs:
        for run in para.runs:
            set_font_on_run(run, cn, en)

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a title at top with underline accent"""
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = DARK
    # underline bar
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.05), Inches(0.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = DARK
    line.line.fill.background()
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(11.9), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(12)
        p2.font.color.rgb = RGBColor(0x99, 0x99, 0x99)

def add_section_label(slide, label):
    txBox = slide.shapes.add_textbox(Inches(10.5), Inches(0.3), Inches(2.5), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    p.alignment = PP_ALIGN.RIGHT

def add_page_number(slide, num, total=21):
    txBox = slide.shapes.add_textbox(Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}" if total else str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0xBB, 0xBB, 0xBB)
    p.alignment = PP_ALIGN.RIGHT

def add_card(slide, left, top, width, height, title, lines, border_color=None):
    """Add a card with title and bullet lines"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = DARK
    p.space_after = Pt(4)
    for line in lines:
        p2 = tf.add_paragraph()
        p2.text = "• " + line
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_GRAY
        p2.space_after = Pt(2)
    set_textframe_lang(tf)
    return shape

def add_table(slide, left, top, width, row_height, headers, rows, col_widths=None):
    """Add a styled table"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(row_height * num_rows))
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x44, 0x44, 0x44)
        set_textframe_lang(cell.text_frame)
    # data
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10.5)
                p.font.color.rgb = TEXT_GRAY
            set_textframe_lang(cell.text_frame)
    return table_shape

def add_text_box(slide, left, top, width, height, text, size=Pt(12), bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    set_textframe_lang(tf)
    return txBox

def add_image_safe(slide, img_name, left, top, width, caption=None):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        pic = slide.shapes.add_picture(path, left, top, width)
        if caption:
            add_text_box(slide, left, top + width * 0.65 + Inches(0.05), width, Inches(0.3),
                         caption, Pt(9), color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)
        return pic
    return None

def add_flow_step(slide, left, top, num, text, step_width=Inches(5.5)):
    """Add a numbered flow step"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.32), Inches(0.32))
    circle.fill.solid()
    circle.fill.fore_color.rgb = DARK
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    set_textframe_lang(tf)
    add_text_box(slide, left + Inches(0.42), top - Pt(1), step_width, Inches(0.35),
                 text, Pt(10.5), color=TEXT_GRAY)

# ============================================================
# SLIDE 1: COVER
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)
# decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(1.8), Inches(0.9), Pt(4))
line.fill.solid(); line.fill.fore_color.rgb = DARK; line.line.fill.background()
add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
             "基于SpringBoot的智能笔记系统的设计与实现",
             Pt(34), bold=True, color=DARK, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.4), Inches(11.3), Inches(0.5),
             "Design and Implementation of an Intelligent Note-Taking System Based on SpringBoot",
             Pt(14), color=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.CENTER)
info_lines = [
    "答辩人：×××",
    "指导教师：×××",
    "所属院系：×××学院",
    "答辩日期：2026年××月"
]
for i, info in enumerate(info_lines):
    add_text_box(slide, Inches(1), Inches(4.3 + i * 0.45), Inches(11.3), Inches(0.4),
                 info, Pt(14), color=TEXT_GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: TOC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, "目  录")
toc_items = [
    ("01", "选题背景", "研究背景 · 国内外现状 · 现存问题 · 研究目标"),
    ("02", "系统设计", "总体架构 · 模块划分 · 技术选型 · 数据库设计 · 安全方案"),
    ("03", "功能实现", "浏览搜索 · 阅读互动 · 智能创作 · AI辅助 · 知识图谱 · 仪表盘 · 测试"),
    ("04", "总  结", "研究成果 · 创新点 · 不足与展望 · 致谢"),
]
for i, (num, title, desc) in enumerate(toc_items):
    y = Inches(1.7 + i * 1.15)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.6), Inches(0.6))
    circle.fill.solid(); circle.fill.fore_color.rgb = DARK; circle.line.fill.background()
    tf = circle.text_frame; p = tf.paragraphs[0]
    p.text = num; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, Inches(1.9), y + Pt(2), Inches(8), Inches(0.35),
                 title, Pt(20), bold=True, color=TEXT_DARK)
    add_text_box(slide, Inches(1.9), y + Inches(0.38), Inches(8), Inches(0.3),
                 desc, Pt(11), color=RGBColor(0x99, 0x99, 0x99))
    # bottom line
    if i < 3:
        bline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.9), y + Inches(0.85), Inches(8.5), Pt(0.5))
        bline.fill.solid(); bline.fill.fore_color.rgb = RGBColor(0xEE, 0xEE, 0xEE); bline.line.fill.background()

# ============================================================
# SLIDE 3: 研究背景
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 01 · 选题背景"); add_page_number(slide, 3)
add_title_bar(slide, "研究背景")
cards_data = [
    ("时代趋势", ["短视频与自媒体广泛流行，用户生成内容爆发式增长，知识的生产和传播方式发生根本性变化。"]),
    ("核心矛盾", ["知识以碎片化形式散落在不同平台和不同时间节点之中——视频笔记、博客收藏、聊天记录、零散草稿各自孤立，难以形成体系。"]),
    ("需求增长", ["随着积累的知识内容日益庞大，对碎片化知识进行系统性整理与结构化管理的需求持续递增，个人知识管理正在成为刚性需求。"]),
    ("技术条件成熟", ["Vue3/SpringBoot等框架高度成熟，大语言模型API逐步开放，ECharts等工具日益完善，为智能化知识管理平台提供了充足的技术支撑。"]),
]
for i, (title, lines) in enumerate(cards_data):
    row = i // 2; col = i % 2
    left = Inches(0.7 + col * 6.0); top = Inches(1.5 + row * 2.8)
    add_card(slide, left, top, Inches(5.6), Inches(2.5), title, lines)

# ============================================================
# SLIDE 4: 国内外研究现状
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 01 · 选题背景"); add_page_number(slide, 4)
add_title_bar(slide, "国内外研究现状")
add_table(slide, Inches(0.7), Inches(1.6), Inches(11.9), 0.9,
          ["类型", "代表产品", "核心特点", "主要不足"],
          [
              ["服务型平台", "Notion、语雀、飞书文档、有道云笔记", "功能丰富、协作便捷、多端同步", "数据存于第三方，存在隐私风险；平台政策变动或停运将导致数据不可控"],
              ["开源框架", "Halo、Typecho、WordPress", "可自行部署、数据掌握在自己手中", "以博客发布为主，缺少笔记管理闭环；缺乏AI辅助写作功能"],
              ["本地笔记工具", "Obsidian、思源笔记、Joplin", "本地优先、Markdown原生支持、插件生态", "侧重单篇编辑体验，知识结构化分析能力不足；缺少数据分析与可视化"],
          ],
          [Inches(1.4), Inches(2.8), Inches(3.0), Inches(4.7)])
add_text_box(slide, Inches(0.7), Inches(5.4), Inches(11.9), Inches(0.8),
             "小结：现有方案在数据自主可控、AI辅助写作和知识结构化分析与可视化三个方向仍存在明显不足，本系统针对这三个缺口进行设计。",
             Pt(12), bold=False, color=TEXT_GRAY)

# ============================================================
# SLIDE 5: 现存问题分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 01 · 选题背景"); add_page_number(slide, 5)
add_title_bar(slide, "现存问题分析")
# left: problems
problems = [
    ("⚠️ 问题一：知识碎片化", "内容散落在视频平台、博客、聊天记录、本地草稿等多个载体中，格式不统一、版本混乱、检索困难。"),
    ("⚠️ 问题二：数据不可控", "主流笔记平台将数据存储在云端服务器，用户无法完全掌控。平台政策调整、服务停运或数据泄露等风险始终存在。"),
    ("⚠️ 问题三：缺乏智能辅助", "现有多数工具仅提供基础编辑与存储能力，缺少AI辅助写作功能和基于数据洞察的写作分析。"),
]
for i, (t, d) in enumerate(problems):
    y = Inches(1.6 + i * 1.7)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, Inches(5.8), Inches(1.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xF2)
    shape.line.color.rgb = RGBColor(0xF5, 0xC6, 0xC6); shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = RGBColor(0xC0, 0x39, 0x2B)
    p2 = tf.add_paragraph(); p2.text = d; p2.font.size = Pt(11); p2.font.color.rgb = TEXT_GRAY
# right: strategy
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.7), Inches(4.5))
shape2.fill.solid(); shape2.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xF0)
shape2.line.color.rgb = RGBColor(0xB8, 0xD8, 0xB8); shape2.line.width = Pt(1)
tf2 = shape2.text_frame; tf2.word_wrap = True; tf2.margin_left = Inches(0.25); tf2.margin_top = Inches(0.15)
p = tf2.paragraphs[0]; p.text = "🎯 本系统的应对策略"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ACCENT_GREEN
strategies = [
    "针对碎片化：统一平台集中管理，通过分类标签、全文搜索和分页浏览实现高效检索",
    "针对数据不可控：全部数据存储于用户自有数据库，采用JWT无状态认证和BCrypt加密保障安全",
    "针对缺乏智能化：集成大语言模型API提供AI辅助能力；通过知识图谱与写作仪表盘将数据结构化呈现",
]
for s in strategies:
    p2 = tf2.add_paragraph(); p2.text = "• " + s; p2.font.size = Pt(11.5); p2.font.color.rgb = TEXT_GRAY; p2.space_after = Pt(8)

# ============================================================
# SLIDE 6: 研究目标与意义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 01 · 选题背景"); add_page_number(slide, 6)
add_title_bar(slide, "研究目标与意义")
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "🎯 研究目标", Pt(16), bold=True, color=DARK)
goals = [
    "覆盖内容创作 → 分类管理 → 阅读互动 → 数据分析的完整闭环",
    "集成大语言模型，提供 AI 辅助写作能力",
    "通过知识图谱将碎片化文章转化为可视化知识网络",
    "通过写作仪表盘实现多维度写作数据统计与分析",
]
for i, g in enumerate(goals):
    y = Inches(2.1 + i * 0.45)
    add_flow_step(slide, Inches(0.9), y, i + 1, g)
add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4), "💡 研究意义", Pt(16), bold=True, color=DARK)
meanings = [
    ("理论意义", "验证了Vue3+SpringBoot+MyBatis技术组合在个人知识管理场景中的可行性与工程化实施路径。"),
    ("实践意义", "为个人创作者提供安全、高效、智能的内容管理平台，解决数据自主可控与知识结构化的痛点。"),
    ("社会意义", "推动个人知识管理工具的智能化发展，促进知识系统化整理与分享，助力终身学习与知识积累。"),
]
for i, (t, d) in enumerate(meanings):
    y = Inches(2.1 + i * 1.5)
    add_card(slide, Inches(7), y, Inches(5.5), Inches(1.3), t, [d])

# ============================================================
# SLIDE 7: 系统总体架构
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 02 · 系统设计"); add_page_number(slide, 7)
add_title_bar(slide, "系统总体架构")
# three layer boxes
layers = [
    ("前端展示层", "Vue 3 + Element Plus  |  Axios · ECharts · wangEditor", RGBColor(0xE3, 0xF2, 0xFD), Inches(5.6), Inches(0.9)),
    ("后端业务层", "Spring Boot 3 + MyBatis  |  JWT · BCrypt · 大模型API", DARK, Inches(5.6), Inches(0.9)),
    ("数据存储层", "MySQL 关系型数据库  |  用户表 · 文章表 · 评论表 · 点赞表 · 收藏表", RGBColor(0xFE, 0xF9, 0xF3), Inches(5.6), Inches(0.9)),
]
for i, (title, desc, color, w, h) in enumerate(layers):
    y = Inches(1.5 + i * 1.6)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.color.rgb = DARK if i == 1 else RGBColor(0xDD, 0xDD, 0xDD); shape.line.width = Pt(1.5)
    tf = shape.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(16); p.font.bold = True
    p.font.color.rgb = WHITE if i == 1 else DARK
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc; p2.font.size = Pt(10.5)
    p2.font.color.rgb = RGBColor(0xEE, 0xEE, 0xFF) if i == 1 else RGBColor(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER
    if i < 2:
        add_text_box(slide, Inches(3.4), y + Inches(0.9), Inches(1), Inches(0.3),
                     "↕", Pt(18), color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.CENTER)
# right side: principles
principles = [
    ("分层解耦", "展示层、业务层、持久层各司其职，通过接口通信降低耦合"),
    ("前后端分离", "前端专注界面交互，后端专注业务逻辑，独立开发部署"),
    ("无状态认证", "JWT令牌机制，服务端不存储会话，支持水平扩展"),
    ("数据安全", "BCrypt加盐哈希存储密码，拦截器验证令牌防未授权访问"),
]
for i, (t, d) in enumerate(principles):
    add_card(slide, Inches(7), Inches(1.5 + i * 1.25), Inches(5.5), Inches(1.1), t, [d])

# ============================================================
# SLIDE 8: 功能模块划分
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 02 · 系统设计"); add_page_number(slide, 8)
add_title_bar(slide, "功能模块划分")
modules = [
    ("❶ 用户认证", "邮箱注册 / 账号登录 / JWT令牌校验\n密码BCrypt加密 / 全局登录状态管理"),
    ("❷ 文章浏览与选择", "五大分类标签筛选 / 关键词模糊搜索\n分页展示 / 文章标题跳转"),
    ("❸ 文章阅读与互动", "Markdown渲染 / 评论发布删除\n点赞收藏切换 / 多格式文件下载"),
    ("❹ 智能创作", "富文本编辑器 / AI续写总结标题优化\n草稿本地暂存 / 在线发布编辑"),
    ("❺ 知识图谱与仪表盘", "关键词提取 / 力导向图可视化\n写作数据多维度统计图表"),
    ("❻ 个人中心", "头像上传 / 信息编辑\n我的文章·收藏·点赞分页管理"),
]
for i, (title, desc) in enumerate(modules):
    row = i // 3; col = i % 3
    left = Inches(0.7 + col * 4.1); top = Inches(1.6 + row * 2.7)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(2.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
    shape.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8); shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Inches(0.2); tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = DARK
    for line in desc.split("\n"):
        p2 = tf.add_paragraph(); p2.text = "· " + line; p2.font.size = Pt(10.5); p2.font.color.rgb = TEXT_GRAY

# ============================================================
# SLIDE 9: 技术选型
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 02 · 系统设计"); add_page_number(slide, 9)
add_title_bar(slide, "关键技术选型")
add_table(slide, Inches(0.7), Inches(1.5), Inches(11.9), 0.42,
          ["层次", "技术", "版本", "选型理由"],
          [
              ["前端框架", "Vue 3", "3.x", "组合式API提升逻辑复用性，响应式系统性能优异，社区活跃"],
              ["UI组件库", "Element Plus", "2.x", "成熟的Vue3组件库，提供表格、表单、分页等丰富UI组件"],
              ["可视化", "ECharts", "5.x", "支持力导向图、柱状图、饼图等20+图表类型，配置式开发效率高"],
              ["HTTP客户端", "Axios", "1.x", "基于Promise的HTTP客户端，支持请求拦截与响应拦截"],
              ["后端框架", "Spring Boot 3", "3.x", "简化Spring应用搭建，内嵌Tomcat，自动配置，开发效率高"],
              ["持久层", "MyBatis", "3.x", "灵活的SQL映射框架，支持注解和XML两种方式，对复杂查询友好"],
              ["数据库", "MySQL", "8.x", "成熟的关系型数据库，事务支持完善，社区资源丰富"],
              ["编辑器", "wangEditor", "5.x", "开源Web富文本编辑器，轻量级、易集成"],
          ],
          [Inches(1.5), Inches(1.8), Inches(1.2), Inches(7.4)])

# ============================================================
# SLIDE 10: 数据库设计
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 02 · 系统设计"); add_page_number(slide, 10)
add_title_bar(slide, "数据库设计")
add_table(slide, Inches(0.7), Inches(1.5), Inches(7.5), 0.55,
          ["数据表", "核心字段", "说明"],
          [
              ["user (用户表)", "id, name, email, password, phone, avatar", "密码经BCrypt加密存储"],
              ["article (文章表)", "id, article_title, article_sign, article_content, author_id, likes, favorites", "正文以Markdown源文本存储"],
              ["comment (评论表)", "id, article_id, user_id, content, create_time", "关联文章与用户信息"],
              ["likes (点赞表)", "id, user_id, article_id", "用户-文章点赞关联记录"],
              ["favorites (收藏表)", "id, user_id, article_id", "用户-文章收藏关联记录"],
          ],
          [Inches(2.2), Inches(3.8), Inches(1.5)])
# right side design points
design_pts = [
    ("关联表模式", "点赞与收藏采用独立关联表，添加时INSERT，取消时DELETE，同步更新文章表计数字段。"),
    ("Markdown源文本存储", "正文以Markdown原始格式存入数据库，前端预览组件实时渲染，下载时按需转换格式。"),
    ("分页查询设计", "通过偏移量计算实现手动分页，每次请求携带页码，后端计算LIMIT起始位置并执行COUNT查询。"),
]
for i, (t, d) in enumerate(design_pts):
    add_card(slide, Inches(8.7), Inches(1.5 + i * 1.9), Inches(4.2), Inches(1.7), t, [d])

# ============================================================
# SLIDE 11: 安全方案与业务流程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 02 · 系统设计"); add_page_number(slide, 11)
add_title_bar(slide, "安全方案与业务流程")
sec_items = [
    ("JWT 无状态认证", "登录后生成JWT令牌（含用户标识、有效期1小时），前端存入localStorage，后续请求在请求头中携带令牌。"),
    ("BCrypt 密码加密", "注册和登录时密码经BCrypt加盐哈希处理，比对时调用验证方法而非明文比较。"),
    ("拦截器白名单", "注册、登录、文章查询等公开接口直接放行；创作、点赞、收藏、下载需校验令牌，失败返回401。"),
]
for i, (t, d) in enumerate(sec_items):
    add_card(slide, Inches(0.7), Inches(1.5 + i * 1.3), Inches(5.8), Inches(1.15), t, [d])
# right: flows
flow_items = [
    ("游客流程", "浏览首页 → 按分类/关键词浏览 → 阅读正文 → 查看评论 → 注册/登录获取完整权限"),
    ("注册用户流程", "登录 → 浏览/搜索 → 阅读互动（评论/点赞/收藏/下载）→ AI辅助创作 → 发布 → 查看图谱仪表盘"),
    ("文章编辑流程", "个人中心 → 我的文章 → 编辑 → 自动加载原文 → 修订 → 提交更新 → 返回列表"),
]
for i, (t, d) in enumerate(flow_items):
    add_card(slide, Inches(7), Inches(1.5 + i * 1.8), Inches(5.5), Inches(1.6), t, [d])

# ============================================================
# SLIDES 12-19: Function Implementation
# ============================================================

# SLIDE 12: 用户认证
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 12)
add_title_bar(slide, "用户认证模块")
auth_cards = [
    ("📝 用户注册", ["前端实时格式校验：邮箱正则匹配、手机号11位验证", "提交后后端检查邮箱唯一性防止重复", "密码经BCrypt加盐哈希后写入数据库", "记录注册时间，成功后引导至登录页"]),
    ("🔑 用户登录", ["根据邮箱查询用户，BCrypt比对密码", "验证通过后生成JWT令牌（HS256签名，1小时有效）", "令牌+用户信息返回前端存入localStorage", "更新全局登录状态，导航栏动态切换菜单项"]),
    ("🛡️ 鉴权与拦截", ["白名单路径（注册/登录/查询/评论列表）直接放行", "其余接口从请求头提取Token，验证签名与有效期", "校验失败返回401，前端提示重新登录"]),
    ("👥 角色权限", ["游客：浏览文章、阅读内容、查看评论", "注册用户：全部功能（创作、互动、下载、数据分析）"]),
]
for i, (t, lines) in enumerate(auth_cards):
    row = i // 2; col = i % 2
    add_card(slide, Inches(0.7 + col * 6.2), Inches(1.5 + row * 2.8), Inches(5.8), Inches(2.5), t, lines)

# SLIDE 13: 浏览与搜索 (with image)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 13)
add_title_bar(slide, "文章浏览与搜索模块")
add_card(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(2.5),
         "🏷️ 分类标签筛选", ["五个标签：全部、日常、技术、杂谈、其他", "点击标签切换分类标识，重置页码", "后端根据分类参数在SQL中添加过滤条件", "每页返回10条记录，含总条数供分页使用"])
add_card(slide, Inches(0.7), Inches(4.3), Inches(5.5), Inches(2.5),
         "🔍 关键词模糊搜索", ["搜索框支持回车或按钮触发", "对文章标题和内容同时LIKE模糊匹配", "搜索结果按时间倒序，支持分页", "搜索模式显示清除按钮，点击恢复分类列表"])
add_image_safe(slide, "浏览界面.png", Inches(6.5), Inches(1.5), Inches(5.8),
               "▲ 文章浏览界面 — 分类标签 + 关键词搜索 + 分页展示")

# SLIDE 14: 阅读与互动 (with images)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 14)
add_title_bar(slide, "文章阅读与互动模块")
add_card(slide, Inches(0.7), Inches(1.5), Inches(3.5), Inches(2.5),
         "📖 文章阅读", ["左右两栏布局：左侧作者信息与作品列表", "右侧使用预览组件实时渲染Markdown", "当前文章高亮，点击其他作品可无缝切换", "同时加载评论、作者作品、点赞收藏状态"])
add_card(slide, Inches(0.7), Inches(4.3), Inches(3.5), Inches(2.5),
         "💬 评论互动", ["输入区+列表展示区两级结构", "点击输入框检查登录状态", "提交后写入评论表并刷新", "本人评论悬浮显示删除按钮"])
add_card(slide, Inches(4.5), Inches(1.5), Inches(3.2), Inches(2.5),
         "❤️ 点赞与收藏", ["切换式操作：首次添加、再次取消", "加载时查询状态初始化按钮样式", "关联表INSERT/DELETE同步更新计数", "已点赞粉色高亮，已收藏金色高亮"])
add_card(slide, Inches(4.5), Inches(4.3), Inches(3.2), Inches(2.5),
         "📥 多格式下载", ["Markdown：原样输出正文", "纯文本TXT：去除Markdown标记", "EPUB：遵循OCF 2.0.1标准", "前端Blob触发浏览器下载含进度提示"])
add_image_safe(slide, "文章详情界面.png", Inches(8), Inches(1.5), Inches(4.5),
               "▲ 文章阅读界面 — Markdown渲染")
add_image_safe(slide, "下载格式弹窗.png", Inches(8), Inches(4.3), Inches(4.5),
               "▲ 多格式下载 — MD / TXT / EPUB")

# SLIDE 15: 智能创作-编辑器 (with image)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 15)
add_title_bar(slide, "智能创作模块——编辑器与草稿")
add_card(slide, Inches(0.7), Inches(1.5), Inches(5.8), Inches(2.5),
         "✍️ 富文本编辑器", ["集成wangEditor开源编辑器，提供丰富的排版工具", "顶部区域：标题输入框 + 分类标签单选组", "支持常用Markdown语法及所见即所得编辑", "右下角悬浮显示保存草稿和发布文章按钮"])
add_card(slide, Inches(0.7), Inches(4.3), Inches(5.8), Inches(2.5),
         "💾 草稿暂存与发布", ["点击保存草稿时将标题、标签、正文序列化为JSON", "存储于浏览器localStorage，下次自动恢复", "发布时校验标题、标签、正文非空", "后端设置创建时间，初始化计数后写入文章表"])
add_image_safe(slide, "创作界面.png", Inches(6.8), Inches(1.5), Inches(5.7),
               "▲ 智能创作界面 — 富文本编辑器 + AI辅助工具 + 草稿暂存与发布")

# SLIDE 16: AI辅助写作
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 16)
add_title_bar(slide, "智能创作模块——AI 辅助写作")
add_card(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(2.5),
         "🤖 大语言模型集成方案", ["前端调用统一后端AI接口，通过动作类型参数区分功能", "后端根据动作类型构建不同的提示词模板", "模板+文章内容+标题打包发送至大模型API", "返回结果封装为响应对象回传前端"])
add_card(slide, Inches(0.7), Inches(4.3), Inches(5.5), Inches(2.5),
         "📋 提示词设计思路", ["续写：请根据以下内容继续撰写...", "总结：请提炼以下文章的核心要点...", "标题优化：请为以下文章生成更有吸引力的标题...", "错误检测：请检查以下文本的错别字和语法问题..."])
ai_funcs = [
    ("📝 续写", "根据上下文继续生成内容，追加到编辑器末尾"),
    ("📊 总结", "自动提炼文章核心观点与要点，追加到末尾"),
    ("💡 标题优化", "生成多条标题建议，以弹窗形式展示"),
    ("🔍 错误检测", "排查错别字与语法问题，弹窗呈现结果"),
]
for i, (t, d) in enumerate(ai_funcs):
    row = i // 2; col = i % 2
    add_card(slide, Inches(6.5 + col * 3.3), Inches(1.5 + row * 1.5), Inches(3.0), Inches(1.3), t, [d])

# SLIDE 17: 知识图谱 (with image)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 17)
add_title_bar(slide, "知识图谱模块")
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.4), "🧠 图谱构建算法", Pt(15), bold=True, color=DARK)
steps = [
    "获取数据：前端请求作者全部文章列表",
    "节点创建：每篇文章生成节点，分类标签生成聚合节点",
    "标题分词：滑动窗口（2~4字）提取词语片段，过滤停用词",
    "关键词匹配：记录每篇文章关键词集合，两两比对",
    "关联连线：共享关键词的文章之间建立连线表示知识关联",
]
for i, s in enumerate(steps):
    add_flow_step(slide, Inches(0.9), Inches(2.1 + i * 0.5), i + 1, s)
add_image_safe(slide, "知识图谱界面.png", Inches(7), Inches(1.5), Inches(5.5),
               "▲ 知识图谱 — 力导向图呈现文章间的关键词关联网络")

# SLIDE 18: 仪表盘 (with image)
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 18)
add_title_bar(slide, "写作数据仪表盘")
add_table(slide, Inches(0.7), Inches(1.5), Inches(6.0), 0.42,
          ["统计指标", "数据来源", "呈现方式"],
          [
              ["文章总数", "COUNT(id)聚合", "渐变色数据卡片"],
              ["总字数", "SUM(LENGTH)聚合", "渐变色数据卡片（含单位换算）"],
              ["收获点赞数", "SUM(likes)聚合", "渐变色数据卡片"],
              ["被收藏数", "SUM(favorites)聚合", "渐变色数据卡片"],
              ["每日发文趋势", "GROUP BY create_time", "柱状图"],
              ["文章分类占比", "GROUP BY article_sign", "环形饼图"],
              ["热门文章TOP5", "ORDER BY likes DESC", "排行列表"],
          ],
          [Inches(1.8), Inches(2.4), Inches(1.8)])
add_image_safe(slide, "仪表盘界面.png", Inches(7), Inches(1.5), Inches(5.5),
               "▲ 写作仪表盘 — 核心数据卡片 + 发文趋势柱状图 + 分类占比饼图")

# SLIDE 19: 个人中心 + 测试
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 03 · 功能实现"); add_page_number(slide, 19)
add_title_bar(slide, "个人中心模块与系统测试")
add_card(slide, Inches(0.7), Inches(1.5), Inches(6.0), Inches(1.5),
         "👤 个人信息管理", ["查看/编辑双模式切换", "头像上传支持JPG/PNG/GIF，不超过2MB", "后端UUID重命名防冲突后保存至本地目录"])
add_card(slide, Inches(0.7), Inches(3.3), Inches(6.0), Inches(1.5),
         "📋 我的文章·收藏·点赞", ["按作者查询文章列表，提供编辑/删除按钮", "我的收藏和点赞通过关联表JOIN查询", "卡片列表展示，点击标题跳转详情，支持分页"])
add_table(slide, Inches(7), Inches(1.5), Inches(5.5), 0.42,
          ["测试模块", "测试项", "结果"],
          [
              ["登录注册", "正确/错误/空值/重复 共9组", "✅ 符合预期"],
              ["首页导航", "5卡片跳转+回退 共10步", "✅ 符合预期"],
              ["文章浏览", "5标签+搜索+分页 共8步", "✅ 符合预期"],
              ["文章详情", "切换/评论/点赞/收藏/下载 共11步", "✅ 符合预期"],
              ["AI写作", "续写/总结/标题/纠错/草稿/发布 共10步", "✅ 符合预期"],
              ["个人中心", "信息编辑/文章管理/收藏点赞 共13步", "✅ 符合预期"],
          ],
          [Inches(1.4), Inches(3.0), Inches(1.1)])

# SLIDE 20: 总结
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide); add_section_label(slide, "PART 04 · 总结"); add_page_number(slide, 20)
add_title_bar(slide, "研究成果与展望")
results = [
    "基于 SpringBoot + Vue3 + MyBatis + MySQL 技术栈完成了个人智能笔记系统的前后端全流程开发",
    "实现了 Markdown 在线编辑、分类检索、评论互动、点赞收藏、多格式下载（MD/TXT/EPUB）等完整功能链",
    "接入大语言模型 API，完成了续写、总结、标题优化与错误检测四项 AI 辅助写作功能",
    "通过标题关键词提取与关联匹配构建知识网络，以力导向图与写作仪表盘实现数据可视化分析",
    "经数十组测试用例验证，各模块运行稳定，核心业务流程均正常完成",
]
add_text_box(slide, Inches(0.7), Inches(1.5), Inches(7), Inches(0.4), "✅ 研究成果", Pt(16), bold=True, color=DARK)
for i, r in enumerate(results):
    add_text_box(slide, Inches(0.9), Inches(2.1 + i * 0.45), Inches(7), Inches(0.4),
                 f"{i+1}. {r}", Pt(11.5), color=TEXT_DARK)
add_text_box(slide, Inches(0.7), Inches(4.8), Inches(7), Inches(0.4),
             "💡 创新点：数据自主可控 · AI辅助写作 · 知识图谱力导向图可视化", Pt(12), bold=True, color=ACCENT_GREEN)

futures = [
    ("关键词提取优化", "当前使用滑动窗口分词匹配，精度有限。后续可引入 TF-IDF 或 NLP 词频算法提升关联质量。"),
    ("AI 功能扩展", "目前支持四项AI功能。下一步可扩展写作风格调整、智能标签推荐、大纲生成等能力。"),
    ("多端适配", "当前仅支持PC端。后续可引入响应式布局，适配平板与手机端，提升移动场景下便捷性。"),
]
add_text_box(slide, Inches(0.7), Inches(5.3), Inches(5), Inches(0.4), "🔮 不足与展望", Pt(16), bold=True, color=DARK)
for i, (t, d) in enumerate(futures):
    add_card(slide, Inches(0.7), Inches(5.8 + i * 0.55), Inches(11.9), Inches(0.48), t, [d])

# SLIDE 21: 致谢
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.2), Inches(2.0), Inches(0.9), Pt(4))
line.fill.solid(); line.fill.fore_color.rgb = DARK; line.line.fill.background()
add_text_box(slide, Inches(1), Inches(2.4), Inches(11.3), Inches(1), "致  谢",
             Pt(36), bold=True, color=DARK, align=PP_ALIGN.CENTER)
thanks = [
    "衷心感谢姜靓老师在本论文选题、方案设计与撰写过程中的悉心指导",
    "感谢各位评审老师在百忙之中对本文的审阅与指正",
    "感谢家人的支持与陪伴，让我能够心无旁骛地完成学业",
    "感谢大学四年间所有授课老师的教导与帮助",
]
for i, t in enumerate(thanks):
    add_text_box(slide, Inches(2), Inches(3.8 + i * 0.5), Inches(9.3), Inches(0.4),
                 t, Pt(15), color=TEXT_GRAY, align=PP_ALIGN.CENTER)
add_text_box(slide, Inches(2), Inches(6.0), Inches(9.3), Inches(0.4),
             "—— 请各位老师批评指正 ——", Pt(13), color=RGBColor(0xBB, 0xBB, 0xBB), align=PP_ALIGN.CENTER)
add_page_number(slide, 21, 21)

# ============================================================
# SAVE
# ============================================================
# Final pass: fix language AND font on ALL text in ALL slides
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            set_textframe_lang(shape.text_frame)
            set_font_on_textframe(shape.text_frame)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    set_textframe_lang(cell.text_frame)
                    set_font_on_textframe(cell.text_frame)

output_path = os.path.join(os.path.dirname(__file__), "毕业论文答辩PPT_v3.pptx")
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
